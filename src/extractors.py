import fitz  # PyMuPDF
from PIL import Image
import pandas as pd
import os
import easyocr
import io
import torch 
from transformers import BlipProcessor, BlipForConditionalGeneration
import zipfile
from lxml import etree
import docx
from pptx import Presentation
from bs4 import BeautifulSoup
import ebooklib
from ebooklib import epub
import re

class LocalExtractor:
    _ocr_reader = None
    _image_captioner = None

    @staticmethod
    def get_ocr_reader():
        """加载 OCR 模型（第一次调用时执行）"""
        if LocalExtractor._ocr_reader is None:
            print("正在加载本地 OCR 模型 (支持中英文)...")
            # 指定支持中英文
            LocalExtractor._ocr_reader = easyocr.Reader(['ch_sim', 'en'], gpu=torch.cuda.is_available())
        return LocalExtractor._ocr_reader

    @staticmethod
    def get_image_captioner():
        """加载图片描述模型 (第一次调用时执行)"""
        if LocalExtractor._image_captioner is None:
            print("正在加载本地视觉模型 (BLIP)...")
            device = "cuda" if torch.cuda.is_available() else "cpu"
            processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
            model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base").to(device)
            LocalExtractor._image_captioner = (processor, model, device)
        return LocalExtractor._image_captioner

    @staticmethod
    def extract_text(file_path):
        """处理纯文本文件"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return (f.read(2000)),"text"  # 只取前2000字作为摘要

    @staticmethod
    def extract_pdf(file_path):
        """处理PDF：提取首页文字和元数据，如果是扫描件则做 OCR"""
        doc = fitz.open(file_path)
        text = ""
        is_scanned = True
        file_type = "pdf"
        
        # 1. 尝试直接提取前 2 页的文字
        for page in doc[:2]:
            text_content = page.get_text()
            if text_content.strip():
                text += text_content
                is_scanned = False
        
        # 2. 如果前2页都没有文字，判定为扫描件，启动本地 OCR
        if is_scanned and len(doc) > 0:
            print(f"检测到扫描版 PDF: {os.path.basename(file_path)}，启动本地 OCR...")
            reader = LocalExtractor.get_ocr_reader()
            
            # 只 OCR 第一页，避免太耗时
            page = doc[0]
            pix = page.get_pixmap()
            img_bytes = pix.tobytes("png")
            
            # 使用 EasyOCR 读取
            result = reader.readtext(img_bytes, detail=0)
            text = " ".join(result)
            file_type = "pdf-ocr"
            
        return (f"[PDF 文本 ({'扫描件' if is_scanned else '电子书'})]: " + text[:2000]), file_type

    @staticmethod
    def extract_image_info(file_path):
        """处理图片：提取EXIF元数据，并使用 BLIP 生成图片描述"""
        print(f"正在智能分析图片: {os.path.basename(file_path)}...")
        
        # A. 提取基础元数据
        try:
            with Image.open(file_path) as img:
                basic_info = f"格式: {img.format}, 尺寸: {img.size}, 模式: {img.mode}"
                # B. 使用 BLIP 生成图片描述
                processor, model, device = LocalExtractor.get_image_captioner()
                # 预处理图片
                inputs = processor(img, return_tensors="pt").to(device)
                # 生成描述
                with torch.no_grad():
                    out = model.generate(**inputs)
                caption = processor.decode(out[0], skip_special_tokens=True)
                return (f"[图片基础信息]: {basic_info}\n[图片AI描述]: {caption}", "image")
                
        except Exception as e:
            return f"[图片解析失败]: {e}"
        
    @staticmethod
    def extract_table(file_path):
        """处理Excel/CSV：提取表头和前几行数据"""
        if file_path.endswith('.csv'):
            df = pd.read_csv(file_path, nrows=5)
        else:
            df = pd.read_excel(file_path, nrows=5)
        return (f"列名: {list(df.columns)}\n数据样例: {df.values.tolist()}"),"table"
    
    @staticmethod
    def extract_word(file_path):
        """处理 Word (.docx) 文件"""
        try:
            doc = docx.Document(file_path)
            # 提取前 10 段文字作为摘要
            full_text = [para.text for para in doc.paragraphs[:10]]
            return "\n".join(full_text)[:2000], "Word文档"
        except Exception as e:
            return f"Word解析失败: {e}", "Word文档(损坏)"

    @staticmethod
    def extract_pptx(file_path):
        """处理 PowerPoint (.pptx) 文件"""
        try:
            prs = Presentation(file_path)
            text_runs = []
            # 遍历前 5 张幻灯片的所有文本框
            for slide in prs.slides[:5]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)[:2000], "PPT演示文稿"
        except Exception as e:
            return f"PPT解析失败: {e}", "PPT(损坏)"

    @staticmethod
    def extract_html(file_path):
        """处理 HTML 文件"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f, 'html.parser')
                # 移除脚本和样式
                for script in soup(["script", "style"]):
                    script.extract()
                return soup.get_text(separator=' ')[:2000], "网页文件"
        except Exception as e:
            return f"HTML解析失败: {e}", "网页文件"

    @staticmethod
    def extract_epub(file_path):
        """
        改进后的 EPUB 提取函数：
        1. 通过 OPF 映射建立逻辑阅读顺序，解决 index_split 分片问题。
        2. 自动处理命名空间，精确定位内容路径。
        3. 200 字符强力兜底策略：排除空格和控制字符。
        """
        content_parts = []
        full_text_for_snippet = ""
        
        try:
            book = epub.read_epub(file_path)
            # 提取书籍元数据：标题和作者
            title = book.get_metadata('DC', 'title')
            creator = clean_epub_authors(book.get_metadata('DC', 'creator'))

            with zipfile.ZipFile(file_path, 'r') as z:
                # --- 1. 定位 OPF 文件 (电子书的地图) ---
                container = z.read('META-INF/container.xml')
                tree = etree.fromstring(container)
                # 获取 content.opf 的路径
                opf_path = tree.xpath('//protocol:rootfile/@full-path', 
                                    namespaces={'protocol': 'urn:oasis:names:tc:opendocument:xmlns:container'})[0]
                base_dir = os.path.dirname(opf_path)

                # --- 2. 解析 OPF 以获取正确的阅读顺序 (Spine) ---
                opf_content = z.read(opf_path)
                opf_tree = etree.fromstring(opf_content)
                ns = {'opf': 'http://www.idpf.org/2007/opf'}
                
                # 建立 ID 到路径的映射 (Manifest)
                manifest = {}
                for item in opf_tree.xpath('//opf:item', namespaces=ns):
                    manifest[item.get('id')] = item.get('href')

                # 按 Spine 定义的顺序读取文件 (解决自私的基因这种分片问题)
                spine_items = opf_tree.xpath('//opf:itemref', namespaces=ns)
                
                for itemref in spine_items:
                    idref = itemref.get('idref')
                    if idref in manifest:
                        html_href = manifest[idref]
                        # 转换路径以适配 zip 内部结构
                        html_path = os.path.join(base_dir, html_href).replace('\\', '/')
                        
                        try:
                            html_content = z.read(html_path)
                            soup = BeautifulSoup(html_content, 'html.parser')
                            
                            # 提取该分片的正文
                            text = soup.get_text(separator=' ', strip=True)
                            if text:
                                content_parts.append(text)
                        except KeyError:
                            continue

                full_text_for_snippet = "".join(content_parts)

        except Exception as e:
            print(f"解析异常: {e}")

        # --- 3. 强力兜底策略 ---
        # 如果通过标准逻辑提取的内容过少（例如少于 50 字），执行字符级暴力提取
        if len(full_text_for_snippet.strip()) < 50:
            full_text_for_snippet = brute_force_extract(file_path)

        # 最终截取前 200 个有效字符（排除空格、换行、控制字符）
        clean_snippet = re.sub(r'[\s\x00-\x1f\x7f-\x9f]', '', full_text_for_snippet)
        summary = f"标题: {title}, 作者: {creator}\n内容片段: {clean_snippet[:200]}"

        return summary, "EPUB电子书"

    
        """处理 EPUB 电子书"""
        try:
            book = epub.read_epub(file_path)
            # 提取书籍元数据：标题和作者
            title = book.get_metadata('DC', 'title')
            creator = book.get_metadata('DC', 'creator')
            
            # 提取第一章节的部分文字
            chapters = list(book.get_items_of_type(ebooklib.ITEM_DOCUMENT))
            print(f"EPUB章节数量: {len(chapters)}")
            first_content = ""
            if chapters:
                soup = BeautifulSoup(chapters[0].get_content(), 'html.parser')
                first_content = soup.get_text()
                
            summary = f"标题: {title}, 作者: {creator}\n内容片段: {first_content}"
            return summary[:2000], "EPUB电子书"
        except Exception as e:
            return f"EPUB解析失败: {e}", "EPUB电子书"

def clean_epub_authors(metadata_authors):
    """
    输入: [('MrToyy', {...}), ('知乎', {...}), ('MrToyy', {...})]
    输出: "MrToyy, 知乎"
    """
    if not metadata_authors:
        return "未知作者"
    
    # 1. 提取元组中的第一个元素（姓名）
    # 2. strip() 清除空格
    # 3. set() 自动去重
    names = []
    seen = set()
    
    for item in metadata_authors:
        # 兼容性处理：有些库返回的是元组，有些可能是字符串
        name = item[0] if isinstance(item, tuple) else str(item)
        name = name.strip()
        
        if name and name not in seen:
            names.append(name)
            seen.add(name)
            
    return ", ".join(names)

def brute_force_extract(epub_path):
    """
    兜底策略：不理会文件结构，遍历 zip 内所有 HTML 提取前段字符
    """
    raw_collected = ""
    try:
        with zipfile.ZipFile(epub_path, 'r') as z:
            # 获取所有 html/htm 文件，按名称排序排除掉封面（通常是 000 或 cover）
            all_files = sorted([f for f in z.namelist() if f.endswith(('.html', '.htm'))])
            for f in all_files:
                text = BeautifulSoup(z.read(f), 'html.parser').get_text()
                # 简单过滤掉明显的版权声明或好讀标题
                if "好讀" in text or "经典版" in text:
                    continue
                raw_collected += text
                if len(re.sub(r'\s', '', raw_collected)) > 500: # 收集够多就停止
                    break
    except:
        pass
    return raw_collected